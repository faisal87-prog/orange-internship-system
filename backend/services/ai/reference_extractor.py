"""Extract readable text from uploaded program reference materials."""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from django.conf import settings

from apps.programs.models import ProgramReferenceMaterial
from services.ai.exceptions import AIValidationError

EXTRACTABLE_EXTENSIONS = {".txt", ".csv", ".pdf", ".docx", ".pptx"}
# Legacy .doc / binary .ppt are not reliably extractable without OCR/extra tooling.
UNSUPPORTED_BUT_KNOWN = {".doc", ".ppt"}


def _max_chars() -> int:
    return int(getattr(settings, "OPENAI_REFERENCE_MAX_CHARS", 80000))


def _decode_text_bytes(raw: bytes) -> str:
    return raw.decode("utf-8", errors="replace")


def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts).strip()


def _extract_docx(raw: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(raw))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts).strip()


def extract_bytes(filename: str, raw: bytes) -> dict[str, Any]:
    """
    Extract full usable text from file bytes.

    Raises AIValidationError if the file is too large to include fully
    (no silent truncation for AI context).
    """
    name = Path(filename).name
    ext = Path(name).suffix.lower()
    max_chars = _max_chars()

    if ext in UNSUPPORTED_BUT_KNOWN:
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": (
                f"Uploaded file '{name}' uses format '{ext}', which cannot be "
                "text-extracted in this MVP. The reference was not analyzed."
            ),
        }

    if ext not in EXTRACTABLE_EXTENSIONS:
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": (
                f"Uploaded file '{name}' is attached but its content type is not "
                "text-extractable for AI roadmap generation."
            ),
        }

    try:
        if ext in {".txt", ".csv"}:
            text = _decode_text_bytes(raw)
        elif ext == ".pdf":
            text = _extract_pdf(raw)
        elif ext == ".docx":
            text = _extract_docx(raw)
        elif ext == ".pptx":
            text = _extract_pptx(raw)
        else:
            text = ""
    except Exception:  # noqa: BLE001
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": (
                f"Uploaded reference '{name}' could not be extracted. "
                "The reference was not analyzed."
            ),
        }

    text = (text or "").strip()
    if not text:
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": (
                f"Uploaded reference '{name}' contained no extractable text. "
                "The reference was not analyzed."
            ),
        }

    if len(text) > max_chars:
        raise AIValidationError(
            f"Reference '{name}' is too large to include fully in AI generation "
            f"(limit {max_chars} characters). Please upload a smaller reference "
            "or split the document, then build the prompt again."
        )

    return {
        "content_retrieved": True,
        "extracted_text": text,
        "content_note": "Full extractable text included for AI roadmap generation.",
    }


def extract_reference_material(material: ProgramReferenceMaterial) -> dict[str, Any]:
    """Extract content for one ProgramReferenceMaterial row."""
    if not material.file:
        if material.external_url:
            return {
                "content_retrieved": False,
                "extracted_text": None,
                "content_note": (
                    "External URL supplied; webpage content was not extracted."
                ),
            }
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": "No file or usable content was available for this reference.",
        }

    name = Path(material.file.name).name
    try:
        with material.file.open("rb") as handle:
            raw = handle.read()
    except Exception:  # noqa: BLE001
        return {
            "content_retrieved": False,
            "extracted_text": None,
            "content_note": (
                f"Uploaded reference '{name}' could not be read. "
                "The reference was not analyzed."
            ),
        }
    return extract_bytes(name, raw)
